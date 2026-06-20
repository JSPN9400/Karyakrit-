"""
Presentation Generator Module

Production-grade AI-powered presentation generation for office use.
"""

import logging
import os
import re
import time
from dataclasses import dataclass
from typing import List, Dict, Optional

from core.llm_provider import LLMProviderManager, PresentationOutline
from pptx import Presentation

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@dataclass
class PresentationRequest:
    """Request data for presentation generation."""
    topic: str
    slide_count: int = 5
    provider: Optional[str] = None  # 'openai', 'gemini', 'ollama', or None for fallback


@dataclass
class SlideContent:
    """Content for a single slide."""
    title: str
    bullets: List[str]


class PresentationGenerator:
    """Production-grade presentation generator with AI and fallback support."""

    def __init__(self):
        """Initialize the generator with configuration."""
        # Ensure output path is always relative to the project root, not current working dir.
        project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
        self.output_base = os.path.join(project_root, 'output', 'presentations')
        os.makedirs(self.output_base, exist_ok=True)

        # Configuration from environment
        self.default_provider = os.getenv('KARYAKRIT_AI_PROVIDER', 'fallback')
        self.openai_api_key = os.getenv('OPENAI_API_KEY')
        self.gemini_api_key = os.getenv('GEMINI_API_KEY')
        self.ollama_url = os.getenv('OLLAMA_URL', 'http://localhost:11434')

        logger.info("PresentationGenerator initialized")

    def generate_presentation(self, request: PresentationRequest) -> str:
        """
        Generate a presentation based on the request.

        Args:
            request: PresentationRequest with topic and parameters.

        Returns:
            str: Path to the generated presentation file.
        """
        logger.info(f"Generating presentation for topic: {request.topic}")

        # Normalize topic
        request.topic = self._normalize_topic(request.topic)

        # Validate request
        if not self._validate_request(request):
            raise ValueError("Invalid presentation request")

        # Generate content
        slides = self._generate_content(request)

        # Validate slides
        if not self._validate_slides(slides):
            raise ValueError("Generated slides are invalid")

        # Create PPT
        output_path = self._generate_output_path(request.topic)
        self._create_pptx(slides, output_path)

        logger.info(f"Presentation saved to: {output_path}")
        return output_path

    def _validate_request(self, request: PresentationRequest) -> bool:
        """Validate the presentation request."""
        if not request.topic or not request.topic.strip():
            logger.error("Topic is empty")
            return False

        if not (3 <= request.slide_count <= 15):
            logger.error(f"Slide count {request.slide_count} out of range (3-15)")
            return False

        return True

    def _generate_content(self, request: PresentationRequest) -> List[SlideContent]:
        """Generate slide content using LLM provider manager."""
        normalized_topic = self._normalize_topic(request.topic)
        manager = LLMProviderManager()
        try:
            outline = manager.generate_presentation_outline(normalized_topic, request.slide_count)
            return self._convert_outline_to_slides(outline)
        except Exception as e:
            logger.warning(f"LLM generation failed: {e}, using local fallback")
            return self._fallback_generate(normalized_topic, request.slide_count)

    def _convert_outline_to_slides(self, outline: PresentationOutline) -> List[SlideContent]:
        """Convert LLM-generated outline to SlideContent format."""
        slides = []
        for slide_data in outline.slides:
            slide = SlideContent(
                title=slide_data.get('title', 'Untitled Slide'),
                bullets=slide_data.get('content', [])
            )
            slides.append(slide)
        return slides

    def _normalize_topic(self, topic: str) -> str:
        """Normalize topic by removing awkward prefixes."""
        topic = topic.strip().lower()
        prefixes = ['on', 'for', 'about', 'regarding', 'concerning']
        words = topic.split()
        if words and words[0] in prefixes:
            return ' '.join(words[1:])
        return topic

    def _fallback_generate(self, topic: str, slide_count: int) -> List[SlideContent]:
        """Generate structured fallback content for office presentations."""
        logger.info(f"Generating fallback content for topic: {topic}")

        # Standard office presentation structure
        structure = [
            {
                "title": f"Introduction to {topic.title()}",
                "bullets": [
                    f"Overview of {topic}",
                    "Importance and relevance",
                    "Objectives of this presentation"
                ]
            },
            {
                "title": f"Current Situation - {topic.title()}",
                "bullets": [
                    "Background and context",
                    "Current trends and developments",
                    "Key statistics and data"
                ]
            },
            {
                "title": f"Key Issues and Analysis - {topic.title()}",
                "bullets": [
                    "Main challenges and opportunities",
                    "Analysis of current approaches",
                    "Critical factors to consider"
                ]
            },
            {
                "title": f"Recommendations - {topic.title()}",
                "bullets": [
                    "Proposed solutions and strategies",
                    "Implementation steps",
                    "Expected outcomes and benefits"
                ]
            },
            {
                "title": f"Risks and Challenges - {topic.title()}",
                "bullets": [
                    "Potential obstacles",
                    "Mitigation strategies",
                    "Contingency plans"
                ]
            },
            {
                "title": f"Conclusion and Next Steps - {topic.title()}",
                "bullets": [
                    "Summary of key points",
                    "Recommendations for action",
                    "Future outlook and monitoring"
                ]
            }
        ]

        # Return requested number of slides
        return [
            SlideContent(
                title=slide["title"],
                bullets=slide["bullets"]
            )
            for slide in structure[:slide_count]
        ]

    def _validate_slides(self, slides: List[SlideContent]) -> bool:
        """Validate that slides have proper content."""
        if not slides:
            return False

        for slide in slides:
            if not slide.title or not slide.bullets:
                return False
            if len(slide.bullets) == 0:
                return False

        return True

    def _create_pptx(self, slides: List[SlideContent], output_path: str):
        """Create PowerPoint presentation from slide content."""
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = slides[0].title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Generated by Karyakrit"

        # Content slides
        bullet_slide_layout = prs.slide_layouts[1]
        for slide_content in slides:
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = slide_content.title

            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""

            for bullet in slide_content.bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        prs.save(output_path)
        logger.info(f"PPT created with {len(slides)} slides")

    def _generate_output_path(self, topic: str) -> str:
        """Generate safe output path for the presentation."""
        safe_topic = self._safe_filename(topic)
        timestamp = int(time.time())
        filename = f"presentation_{safe_topic}_{timestamp}.pptx"
        return os.path.join(self.output_base, filename)

    def _safe_filename(self, text: str) -> str:
        """Create a safe filename from text."""
        # Remove special characters and replace spaces with underscores
        safe = re.sub(r'[^\w\s-]', '', text)
        safe = re.sub(r'[-\s]+', '_', safe)
        return safe.lower()[:50]  # Limit length